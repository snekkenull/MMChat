import torch
from PIL import Image
import gradio as gr
from transformers import AutoModelForCausalLM, AutoTokenizer, TextIteratorStreamer
import os
from threading import Thread

import pymupdf
import docx
from pptx import Presentation


MODEL_LIST = ["THUDM/glm-4v-9b"]
HF_TOKEN = os.environ.get("HF_TOKEN", None)
MODEL_ID = os.environ.get("MODEL_ID")
MODEL_NAME = MODEL_ID.split("/")[-1]

TITLE = "<h1>VL-Chatbox</h1>"

DESCRIPTION = f"""
<center>
<p>😊 A Space For My Fav Multimodal.
<br>
🚀 MODEL NOW: <a href="https://hf.co/{MODEL_ID}">{MODEL_NAME}</a>
<br>
✨ Tips: Now you can send DM or upload 1 IMAGE/FILE per time.
<br>
✨ Tips: Please increase MAX LENGTH when deal with file.
<br>
🤙 Supported Format: pdf, txt, docx, pptx, md, png, jpg, webp
<br>
🙇‍♂️ May be rebuilding from time to time.</p>
</center>"""

CSS = """
h1 {
    text-align: center;
    display: block;
}
"""

model = AutoModelForCausalLM.from_pretrained(
    MODEL_ID,
    torch_dtype=torch.bfloat16,
    low_cpu_mem_usage=True,
    trust_remote_code=True
).to(0)
tokenizer = AutoTokenizer.from_pretrained(MODEL_ID, trust_remote_code=True)
model.eval()


def extract_text(path):
    return open(path, 'r').read()

def extract_pdf(path):
    doc = pymupdf.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_docx(path):
    doc = docx.Document(path)
    data = []
    for paragraph in doc.paragraphs:
        data.append(paragraph.text)
    content = '\n\n'.join(data)
    return content

def extract_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def mode_load(path):
    choice = ""
    file_type = path.split(".")[-1]
    print(file_type)
    if file_type in ["pdf", "txt", "py", "docx", "pptx", "json", "cpp", "md"]:
        if file_type.endswith("pdf"):
            content = extract_pdf(path)
        elif file_type.endswith("docx"):
            content = extract_docx(path)
        elif file_type.endswith("pptx"):
            content = extract_pptx(path)
        else:
            content = extract_text(path)
        choice = "doc"
        print(content[:100])
        return choice, content[:5000]
    elif file_type in ["png", "jpg", "jpeg", "bmp", "tiff", "webp"]:
        content = Image.open(path).convert('RGB')
        choice = "image"
        return choice, content
    else:
        raise gr.Error("Oops, unsupported files.")

def stream_chat(message, history: list, temperature: float, max_length: int, top_p: float, top_k: int, penalty: float):
    print(f'message is - {message}')
    print(f'history is - {history}')
    conversation = []
    prompt_files = []
    if message["files"]:
        choice, contents = mode_load(message["files"][-1])
        if choice == "image":
            conversation.append({"role": "user", "image": contents, "content": message['text']})
        elif choice == "doc":
            format_msg = contents + "\n\n\n" + "{} files uploaded.\n" + message['text']
            conversation.append({"role": "user", "content": format_msg})
    else:
        if len(history) == 0:
            #raise gr.Error("Please upload an image first.")
            contents = None
            conversation.append({"role": "user", "content": message['text']})
        else:
            #image = Image.open(history[0][0][0])
            for prompt, answer in history:
                if answer is None:
                    prompt_files.append(prompt[0])
                    conversation.extend([{"role": "user", "content": ""},{"role": "assistant", "content": ""}])
                else:
                    conversation.extend([{"role": "user", "content": prompt}, {"role": "assistant", "content": answer}])
            choice, contents = mode_load(prompt_files[-1])
            if choice == "image":
                conversation.append({"role": "user", "image": contents, "content": message['text']})
            elif choice == "doc":
                format_msg = contents + "\n\n\n" + "{} files uploaded.\n" + message['text']
                conversation.append({"role": "user", "content": format_msg})
    print(f"Conversation is -\n{conversation}")

    input_ids = tokenizer.apply_chat_template(conversation, tokenize=True, add_generation_prompt=True, return_tensors="pt", return_dict=True).to(model.device)
    streamer = TextIteratorStreamer(tokenizer, timeout=60.0, skip_prompt=True, skip_special_tokens=True)

    generate_kwargs = dict(
        max_length=max_length,
        streamer=streamer,
        do_sample=True,
        top_p=top_p,
        top_k=top_k,
        temperature=temperature,
        repetition_penalty=penalty,
        eos_token_id=[151329, 151336, 151338],
    )
    gen_kwargs = {**input_ids, **generate_kwargs}

    with torch.no_grad():
        thread = Thread(target=model.generate, kwargs=gen_kwargs)
        thread.start()
        buffer = ""
        for new_text in streamer:
            buffer += new_text
            yield buffer
 



chatbot = gr.Chatbot()
chat_input = gr.MultimodalTextbox(
    interactive=True,
    placeholder="Enter message or upload a file one time...",
    show_label=False,

)
EXAMPLES = [
        [{"text": "Describe it in detailed", "files": ["./laptop.jpg"]}],
        [{"text": "Where it is?", "files": ["./hotel.jpg"]}],
        [{"text": "Is it real?", "files": ["./spacecat.png"]}]
]

with gr.Blocks(css=CSS, theme="soft",fill_height=True) as demo:
    gr.HTML(TITLE)
    gr.HTML(DESCRIPTION)
    gr.ChatInterface(
        fn=stream_chat,
        multimodal=True,
        textbox=chat_input,
        chatbot=chatbot,
        fill_height=True,
        additional_inputs_accordion=gr.Accordion(label="⚙️ Parameters", open=False, render=False),
        additional_inputs=[
            gr.Slider(
                minimum=0,
                maximum=1,
                step=0.1,
                value=0.8,
                label="Temperature",
                render=False,
            ),
            gr.Slider(
                minimum=1024,
                maximum=8192,
                step=1,
                value=4096,
                label="Max Length",
                render=False,
            ),
            gr.Slider(
                minimum=0.0,
                maximum=1.0,
                step=0.1,
                value=1.0,
                label="top_p",
                render=False,
            ),
            gr.Slider(
                minimum=1,
                maximum=20,
                step=1,
                value=10,
                label="top_k",
                render=False,
            ),
            gr.Slider(
                minimum=0.0,
                maximum=2.0,
                step=0.1,
                value=1.0,
                label="Repetition penalty",
                render=False,
            ),
        ],
    ),
    gr.Examples(EXAMPLES,[chat_input])


if __name__ == "__main__":
    demo.queue(api_open=False).launch(show_api=False, share=False)